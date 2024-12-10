import os
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from pptx import Presentation
from pptx.util import Inches
import itertools
import requests
import urllib3
from PIL import Image
import io
import logging
import base64

# Suppress SSL warnings
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

class TableauRESTAPIProcessor:
    def __init__(self, server_url, site_id, username, password):
        """
        Initialize Tableau REST API connection parameters
        
        :param server_url: Tableau Server URL
        :param site_id: Tableau Site ID
        :param username: Tableau Server username
        :param password: Tableau Server password
        """
        # Logging setup
        logging.basicConfig(level=logging.INFO, 
                            format='%(asctime)s - %(levelname)s: %(message)s')
        self.logger = logging.getLogger(__name__)
        
        # Connection parameters
        self.server_url = server_url.rstrip('/')
        self.site_id = site_id
        self.username = username
        self.password = password
        
        # API endpoint URLs
        self.base_url = f"{self.server_url}/api/3.21"
        
        # Authentication tokens
        self.site_token = None
        self.user_token = None
        
        # Workbook and filter details
        self.workbook_id = None
        self.view_id = None
        self.filters = {}
        
        # Main application window
        self.root = None
        
        # Filter selection variables
        self.filter_vars = {}
    
    def authenticate(self):
        """
        Authenticate with Tableau REST API and get site and user tokens
        """
        try:
            # Prepare authentication request
            auth_url = f"{self.base_url}/auth/signin"
            payload = {
                "credentials": {
                    "name": self.username,
                    "password": self.password,
                    "site": {
                        "contentUrl": self.site_id
                    }
                }
            }
            
            # Send authentication request
            response = requests.post(
                auth_url, 
                json=payload, 
                verify=False  # Use only in controlled environments
            )
            
            # Check authentication response
            response.raise_for_status()
            auth_response = response.json()
            
            # Extract tokens
            self.site_token = auth_response['credentials']['token']
            self.user_token = auth_response['credentials']['user']['id']
            
            self.logger.info("Successfully authenticated with Tableau Server")
            return True
        
        except requests.exceptions.RequestException as e:
            self.logger.error(f"Authentication failed: {e}")
            messagebox.showerror("Authentication Error", str(e))
            return False
    
    def get_workbooks(self):
        """
        Retrieve list of workbooks for the authenticated user
        """
        try:
            # Prepare headers
            headers = {
                'X-Tableau-Auth': self.site_token
            }
            
            # Get workbooks endpoint
            workbooks_url = f"{self.base_url}/sites/{self.site_id}/workbooks"
            
            # Send request
            response = requests.get(
                workbooks_url, 
                headers=headers, 
                verify=False
            )
            
            response.raise_for_status()
            workbooks = response.json()
            
            return workbooks['workbooks']['workbook']
        
        except requests.exceptions.RequestException as e:
            self.logger.error(f"Failed to retrieve workbooks: {e}")
            messagebox.showerror("Workbook Retrieval Error", str(e))
            return []
    
    def get_workbook_details(self, workbook_id):
        """
        Get details of a specific workbook
        
        :param workbook_id: ID of the workbook
        :return: Workbook details
        """
        try:
            # Prepare headers
            headers = {
                'X-Tableau-Auth': self.site_token
            }
            
            # Get workbook details endpoint
            workbook_url = f"{self.base_url}/sites/{self.site_id}/workbooks/{workbook_id}"
            
            # Send request
            response = requests.get(
                workbook_url, 
                headers=headers, 
                verify=False
            )
            
            response.raise_for_status()
            return response.json()
        
        except requests.exceptions.RequestException as e:
            self.logger.error(f"Failed to retrieve workbook details: {e}")
            messagebox.showerror("Workbook Details Error", str(e))
            return None
    
    def get_workbook_filters(self, workbook_id, view_id):
        """
        Retrieve filters for a specific workbook view
        
        :param workbook_id: ID of the workbook
        :param view_id: ID of the view
        :return: Dictionary of filters
        """
        try:
            # Prepare headers
            headers = {
                'X-Tableau-Auth': self.site_token
            }
            
            # Get filters endpoint
            filters_url = (f"{self.base_url}/sites/{self.site_id}/workbooks/"
                           f"{workbook_id}/views/{view_id}/filters")
            
            # Send request
            response = requests.get(
                filters_url, 
                headers=headers, 
                verify=False
            )
            
            response.raise_for_status()
            filter_data = response.json()
            
            # Process filters
            filters = {}
            for filter_item in filter_data.get('filters', {}).get('filter', []):
                # Extract filter name and available values
                filter_name = filter_item.get('name', 'Unknown')
                
                # Get filter values (this might need adjustment based on Tableau's response)
                filter_values = self._get_filter_values(
                    workbook_id, view_id, filter_name
                )
                
                filters[filter_name] = filter_values
            
            return filters
        
        except requests.exceptions.RequestException as e:
            self.logger.error(f"Failed to retrieve filters: {e}")
            messagebox.showerror("Filter Retrieval Error", str(e))
            return {}
    
    def _get_filter_values(self, workbook_id, view_id, filter_name):
        """
        Get available values for a specific filter
        
        :param workbook_id: ID of the workbook
        :param view_id: ID of the view
        :param filter_name: Name of the filter
        :return: List of filter values
        """
        try:
            # Prepare headers
            headers = {
                'X-Tableau-Auth': self.site_token
            }
            
            # Get filter values endpoint
            values_url = (f"{self.base_url}/sites/{self.site_id}/workbooks/"
                          f"{workbook_id}/views/{view_id}/filters/{filter_name}/values")
            
            # Send request
            response = requests.get(
                values_url, 
                headers=headers, 
                verify=False
            )
            
            response.raise_for_status()
            values_data = response.json()
            
            # Extract and return filter values
            return [
                value.get('value', '') 
                for value in values_data.get('filterValues', {}).get('filterValue', [])
            ]
        
        except requests.exceptions.RequestException as e:
            self.logger.error(f"Failed to retrieve filter values: {e}")
            return []
    
    def capture_view_image(self, workbook_id, view_id, filter_combination):
        """
        Capture image of the view with applied filters
        
        :param workbook_id: ID of the workbook
        :param view_id: ID of the view
        :param filter_combination: Dictionary of filter values to apply
        :return: Base64 encoded image
        """
        try:
            # Prepare headers
            headers = {
                'X-Tableau-Auth': self.site_token
            }
            
            # Prepare filter parameters
            filter_params = []
            for filter_name, filter_value in filter_combination.items():
                filter_params.append(
                    f"filter_{filter_name}={urllib3.quote(str(filter_value))}"
                )
            
            # Construct image download URL
            image_url = (f"{self.server_url}/api/3.21/sites/{self.site_id}/workbooks/"
                         f"{workbook_id}/views/{view_id}/image?")
            
            # Add filter parameters if any
            if filter_params:
                image_url += '&'.join(filter_params)
            
            # Send request to get image
            response = requests.get(
                image_url, 
                headers=headers, 
                verify=False
            )
            
            response.raise_for_status()
            
            # Return base64 encoded image
            return base64.b64encode(response.content).decode('utf-8')
        
        except requests.exceptions.RequestException as e:
            self.logger.error(f"Failed to capture view image: {e}")
            return None
    
    def create_ui(self):
        """
        Create user interface for workbook and view selection
        """
        # Create main window
        self.root = tk.Tk()
        self.root.title("Tableau Dashboard Filter Processor")
        self.root.geometry("600x500")
        
        # Workbook selection
        tk.Label(self.root, text="Select Workbook:").pack(pady=(10,0))
        workbook_listbox = tk.Listbox(self.root, width=50)
        workbook_listbox.pack(pady=5)
        
        # Populate workbooks
        workbooks = self.get_workbooks()
        for wb in workbooks:
            workbook_listbox.insert(tk.END, f"{wb['name']} (ID: {wb['id']})")
        
        # View selection
        tk.Label(self.root, text="Select View:").pack(pady=(10,0))
        view_listbox = tk.Listbox(self.root, width=50)
        view_listbox.pack(pady=5)
        
        def on_workbook_select(event):
            # Clear previous view selections
            view_listbox.delete(0, tk.END)
            
            # Get selected workbook
            selection = workbook_listbox.curselection()
            if not selection:
                return
            
            selected_wb = workbooks[selection[0]]
            workbook_id = selected_wb['id']
            
            # Fetch workbook details to get views
            wb_details = self.get_workbook_details(workbook_id)
            
            if wb_details and 'views' in wb_details:
                views = wb_details['views']['view']
                for view in views:
                    view_listbox.insert(tk.END, f"{view['name']} (ID: {view['id']})")
        
        def on_view_select(event):
            # Get selected view
            wb_selection = workbook_listbox.curselection()
            view_selection = view_listbox.curselection()
            
            if not wb_selection or not view_selection:
                return
            
            selected_wb = workbooks[wb_selection[0]]
            selected_view = selected_wb['views']['view'][view_selection[0]]
            
            self.workbook_id = selected_wb['id']
            self.view_id = selected_view['id']
            
            # Fetch filters for the selected view
            self.filters = self.get_workbook_filters(self.workbook_id, self.view_id)
            
            # Close current window and open filter selection
            self.root.destroy()
            self.create_filter_selection_ui()
        
        # Bind selection events
        workbook_listbox.bind('<<ListboxSelect>>', on_workbook_select)
        view_listbox.bind('<<ListboxSelect>>', on_view_select)
        
        self.root.mainloop()
    
    def create_filter_selection_ui(self):
        """
        Create UI for selecting filter combinations
        """
        self.root = tk.Tk()
        self.root.title("Tableau Filter Selection")
        self.root.geometry("600x500")
        
        # Filter selection frames
        for filter_name, options in self.filters.items():
            frame = ttk.LabelFrame(self.root, text=filter_name)
            frame.pack(padx=10, pady=10, fill='x')
            
            # Multiselect listbox
            self.filter_vars[filter_name] = []
            listbox = tk.Listbox(frame, selectmode=tk.MULTIPLE, height=5)
            for option in options:
                listbox.insert(tk.END, option)
            listbox.pack(padx=10, pady=10, fill='x')
            
            # Bind selection
            listbox.bind('<<ListboxSelect>>', 
                         lambda e, name=filter_name: self.update_filter_selection(e, name))
        
        # Generate button
        generate_btn = ttk.Button(self.root, text="Generate PowerPoint", command=self.generate_ppt)
        generate_btn.pack(pady=20)
        
        self.root.mainloop()
    
    def update_filter_selection(self, event, filter_name):
        """
        Update selected filter options
        
        :param event: Tkinter event
        :param filter_name: Name of the filter
        """
        listbox = event.widget
        selected_indices = listbox.curselection()
        
        # Update filter selections
        self.filter_vars[filter_name] = [
            listbox.get(idx) for idx in selected_indices
        ]
    
    def generate_ppt(self):
        """
        Generate PowerPoint with all filter combinations
        """
        # Validate filter selections
        for filter_name, selections in self.filter_vars.items():
            if not selections:
                messagebox.showerror("Error", f"Please select at least one option for {filter_name}")
                return
        
        # Generate filter combinations
        combinations = list(itertools.product(
            *[self.filter_vars[filter_name] for filter_name in self.filters.keys()]
        ))
        
        # Choose save location
        save_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        
        if not save_path:
            return
        
        # Create PowerPoint
        prs = Presentation()
        
        # Add slides for each combination
        for i, combo in enumerate(combinations, 1):
            # Create combination dictionary
            combo_dict = dict(zip(self.filters.keys(), combo))
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            
            # Add text with filter combination details
            title = slide.shapes.title
            title.text = f"Combination {i}"
            
            # Add text box with filter details
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(6), Inches(2)
            )
            tf = txBox.text_frame
            
            # Add filter details to text box
            for j, (filter_name, value) in enumerate(combo_dict.items()):
                p = tf.add_paragraph()
                p.text = f"{filter_name}: {value}"
            
            # Capture
import os
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from pptx import Presentation
from pptx.util import Inches
import itertools
import requests
import urllib3
from PIL import Image
import io
import logging
import base64

# Suppress SSL warnings
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

class TableauRESTAPIProcessor:
    # ... (previous code remains the same)

    def generate_ppt(self):
        """
        Generate PowerPoint with all filter combinations
        """
        # Validate filter selections
        for filter_name, selections in self.filter_vars.items():
            if not selections:
                messagebox.showerror("Error", f"Please select at least one option for {filter_name}")
                return
        
        # Generate filter combinations
        combinations = list(itertools.product(
            *[self.filter_vars[filter_name] for filter_name in self.filters.keys()]
        ))
        
        # Choose save location
        save_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        
        if not save_path:
            return
        
        # Create PowerPoint
        prs = Presentation()
        
        # Add slides for each combination
        for i, combo in enumerate(combinations, 1):
            # Create combination dictionary
            combo_dict = dict(zip(self.filters.keys(), combo))
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            
            # Add text with filter combination details
            title = slide.shapes.title
            title.text = f"Combination {i}"
            
            # Add text box with filter details
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(6), Inches(2)
            )
            tf = txBox.text_frame
            
            # Add filter details to text box
            for j, (filter_name, value) in enumerate(combo_dict.items()):
                p = tf.add_paragraph()
                p.text = f"{filter_name}: {value}"
            
            # Capture view image with applied filters
            try:
                image_base64 = self.capture_view_image(
                    self.workbook_id, 
                    self.view_id, 
                    combo_dict
                )
                
                if image_base64:
                    # Decode base64 image
                    image_data = base64.b64decode(image_base64)
                    
                    # Save temporary image
                    with open('temp_tableau_image.png', 'wb') as f:
                        f.write(image_data)
                    
                    # Add image to slide
                    slide.shapes.add_picture(
                        'temp_tableau_image.png', 
                        Inches(1), Inches(4), 
                        width=Inches(6)
                    )
                    
                    # Remove temporary image
                    os.remove('temp_tableau_image.png')
            
            except Exception as e:
                self.logger.error(f"Error adding image to slide: {e}")
                messagebox.showwarning("Image Capture", f"Could not capture image for combination {i}")
        
        # Save PowerPoint
        prs.save(save_path)
        messagebox.showinfo("Success", f"PowerPoint generated with {len(combinations)} slides")


def main():
    """
    Main entry point for the Tableau Dashboard Filter Processor
    """
    # Create login window
    login_window = tk.Tk()
    login_window.title("Tableau Dashboard Processor")
    login_window.geometry("400x400")

    # Server URL Label and Entry
    tk.Label(login_window, text="Tableau Server URL:").pack(pady=(10,0))
    server_url_entry = tk.Entry(login_window, width=50)
    server_url_entry.pack(pady=5)
    server_url_entry.insert(0, "https://your-tableau-server.com")

    # Site ID Label and Entry
    tk.Label(login_window, text="Site ID:").pack(pady=(10,0))
    site_id_entry = tk.Entry(login_window, width=50)
    site_id_entry.pack(pady=5)

    # Username Label and Entry
    tk.Label(login_window, text="Username:").pack(pady=(10,0))
    username_entry = tk.Entry(login_window, width=50)
    username_entry.pack(pady=5)

    # Password Label and Entry
    tk.Label(login_window, text="Password:").pack(pady=(10,0))
    password_entry = tk.Entry(login_window, show="*", width=50)
    password_entry.pack(pady=5)

    def start_processing():
        """
        Initialize the Tableau REST API Processor and start the workflow
        """
        # Validate inputs
        if not all([
            server_url_entry.get(), 
            site_id_entry.get(), 
            username_entry.get(), 
            password_entry.get()
        ]):
            messagebox.showerror("Error", "Please fill in all fields")
            return

        # Create processor instance
        processor = TableauRESTAPIProcessor(
            server_url_entry.get(),
            site_id_entry.get(),
            username_entry.get(),
            password_entry.get()
        )

        # Authenticate
        if processor.authenticate():
            # Close login window
            login_window.destroy()
            
            # Start UI for workbook and view selection
            processor.create_ui()

    # Login Button
    login_btn = tk.Button(login_window, text="Connect", command=start_processing)
    login_btn.pack(pady=20)

    # Start the login window
    login_window.mainloop()

if __name__ == "__main__":
    main()
