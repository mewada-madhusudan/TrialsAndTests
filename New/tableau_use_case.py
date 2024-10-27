import sys
import json
import os
from datetime import datetime
from io import BytesIO
from typing import List, Dict

from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QFrame, QScrollArea,
    QGroupBox, QGridLayout, QMessageBox, QProgressDialog
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
from PyQt6.QtGui import QIcon, QFont
import tableauserverclient as TSC
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import itertools


class FilterRow(QFrame):
    deleted = pyqtSignal(object)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setup_ui()

    def setup_ui(self):
        layout = QHBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        # Filter name input
        self.name_input = QLineEdit()
        self.name_input.setPlaceholderText("Filter Name")
        self.name_input.setMinimumWidth(150)

        # Field name input
        self.field_input = QLineEdit()
        self.field_input.setPlaceholderText("Field Name")
        self.field_input.setMinimumWidth(150)

        # Options input
        self.options_input = QLineEdit()
        self.options_input.setPlaceholderText("Options (comma-separated)")
        self.options_input.setMinimumWidth(200)

        # Delete button
        delete_btn = QPushButton("×")
        delete_btn.setFixedSize(30, 30)
        delete_btn.clicked.connect(lambda: self.deleted.emit(self))

        # Add widgets to layout
        layout.addWidget(QLabel("Name:"))
        layout.addWidget(self.name_input)
        layout.addWidget(QLabel("Field:"))
        layout.addWidget(self.field_input)
        layout.addWidget(QLabel("Options:"))
        layout.addWidget(self.options_input)
        layout.addWidget(delete_btn)

    def get_data(self) -> Dict:
        return {
            'name': self.name_input.text(),
            'field': self.field_input.text(),
            'options': [opt.strip() for opt in self.options_input.text().split(',') if opt.strip()]
        }


class GeneratorThread(QThread):
    progress = pyqtSignal(int)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, config, view_url, filters):
        super().__init__()
        self.config = config
        self.view_url = view_url
        self.filters = filters

    def run(self):
        try:
            generator = TableauScreenshotGenerator(self.config, self.view_url, self.filters)
            generator.progress_signal = self.progress
            filename = generator.generate_presentation()
            self.finished.emit(filename)
        except Exception as e:
            self.error.emit(str(e))


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.filter_rows: List[FilterRow] = []
        self.setup_ui()

    def setup_ui(self):
        self.setWindowTitle("Tableau Dashboard Screenshot Generator")
        self.setMinimumSize(800, 600)

        # Create central widget and main layout
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)

        # Create configuration group
        config_group = QGroupBox("Tableau Configuration")
        config_layout = QGridLayout()

        # Server configuration fields
        self.server_url = QLineEdit()
        self.server_url.setPlaceholderText("https://your-tableau-server")
        self.site_id = QLineEdit()
        self.site_id.setPlaceholderText("your-site-id")
        self.token_name = QLineEdit()
        self.token_name.setPlaceholderText("your-token-name")
        self.token_value = QLineEdit()
        self.token_value.setPlaceholderText("your-personal-access-token")
        self.token_value.setEchoMode(QLineEdit.EchoMode.Password)

        # Add configuration fields to layout
        config_layout.addWidget(QLabel("Server URL:"), 0, 0)
        config_layout.addWidget(self.server_url, 0, 1)
        config_layout.addWidget(QLabel("Site ID:"), 1, 0)
        config_layout.addWidget(self.site_id, 1, 1)
        config_layout.addWidget(QLabel("Token Name:"), 2, 0)
        config_layout.addWidget(self.token_name, 2, 1)
        config_layout.addWidget(QLabel("Token Value:"), 3, 0)
        config_layout.addWidget(self.token_value, 3, 1)

        config_group.setLayout(config_layout)
        main_layout.addWidget(config_group)

        # View URL group
        url_group = QGroupBox("Dashboard View URL")
        url_layout = QHBoxLayout()
        self.view_url = QLineEdit()
        self.view_url.setPlaceholderText("Enter the Tableau dashboard view URL")
        url_layout.addWidget(self.view_url)
        url_group.setLayout(url_layout)
        main_layout.addWidget(url_group)

        # Filters group
        filters_group = QGroupBox("Filters Configuration")
        filters_layout = QVBoxLayout()

        # Scroll area for filters
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll_widget = QWidget()
        self.filters_layout = QVBoxLayout(scroll_widget)
        scroll.setWidget(scroll_widget)

        # Add filter button
        add_filter_btn = QPushButton("Add Filter")
        add_filter_btn.clicked.connect(self.add_filter_row)
        filters_layout.addWidget(add_filter_btn)
        filters_layout.addWidget(scroll)

        filters_group.setLayout(filters_layout)
        main_layout.addWidget(filters_group)

        # Generate button
        generate_btn = QPushButton("Generate PPT")
        generate_btn.setStyleSheet("background-color: #4CAF50; color: white; padding: 10px;")
        generate_btn.clicked.connect(self.generate_screenshots)
        main_layout.addWidget(generate_btn)

    def add_filter_row(self):
        filter_row = FilterRow()
        filter_row.deleted.connect(self.remove_filter_row)
        self.filters_layout.addWidget(filter_row)
        self.filter_rows.append(filter_row)

    def remove_filter_row(self, row):
        self.filters_layout.removeWidget(row)
        self.filter_rows.remove(row)
        row.deleteLater()

    def get_config(self):
        return {
            'server_url': self.server_url.text(),
            'site_id': self.site_id.text(),
            'token_name': self.token_name.text(),
            'token_value': self.token_value.text()
        }

    def generate_screenshots(self):
        # Validate inputs
        if not all([self.server_url.text(), self.site_id.text(),
                    self.token_name.text(), self.token_value.text()]):
            QMessageBox.warning(self, "Validation Error", "Please fill in all Tableau configuration fields.")
            return

        if not self.view_url.text():
            QMessageBox.warning(self, "Validation Error", "Please enter a dashboard view URL.")
            return

        if not self.filter_rows:
            QMessageBox.warning(self, "Validation Error", "Please add at least one filter configuration.")
            return

        # Get filter configurations
        filters = [row.get_data() for row in self.filter_rows]

        # Create and start generator thread
        self.generator_thread = GeneratorThread(
            self.get_config(),
            self.view_url.text(),
            filters
        )

        # Create progress dialog
        self.progress_dialog = QProgressDialog("Generating screenshots...", "Cancel", 0, 100, self)
        self.progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)

        # Connect signals
        self.generator_thread.progress.connect(self.progress_dialog.setValue)
        self.generator_thread.finished.connect(self.generation_finished)
        self.generator_thread.error.connect(self.generation_error)

        # Start generation
        self.generator_thread.start()

    def generation_finished(self, filename):
        self.progress_dialog.close()
        QMessageBox.information(self, "Success",
                                f"PowerPoint presentation has been generated successfully!\nSaved as: {filename}")

    def generation_error(self, error_message):
        self.progress_dialog.close()
        QMessageBox.critical(self, "Error", f"An error occurred: {error_message}")


class TableauScreenshotGenerator:
    def __init__(self, config, view_url, filter_configs):
        self.config = config
        self.view_url = view_url
        self.filter_configs = filter_configs
        self.tableau_auth = None
        self.server = None
        self.progress_signal = None

    def setup_tableau_client(self):
        self.tableau_auth = TSC.PersonalAccessTokenAuth(
            token_name=self.config['token_name'],
            personal_access_token=self.config['token_value'],
            site_id=self.config['site_id']
        )
        self.server = TSC.Server(self.config['server_url'])

    def get_view_image(self, filter_combination):
        filter_params = {}
        for filter_config, value in zip(self.filter_configs, filter_combination):
            filter_params[f"vf_{filter_config['field']}"] = value

        self.server.auth.sign_in(self.tableau_auth)
        view_id = self.extract_view_id_from_url(self.view_url)
        view = self.server.views.get_by_id(view_id)
        image_response = self.server.views.get_image(view, filter_params)

        return Image.open(BytesIO(image_response))

    def extract_view_id_from_url(self, url):
        # Implement based on your URL format
        return url.split('/')[-1]

    def generate_presentation(self) -> str:
        try:
            self.setup_tableau_client()
            prs = Presentation()

            # Calculate total combinations
            filter_options = [config['options'] for config in self.filter_configs]
            combinations = list(itertools.product(*filter_options))
            total_combinations = len(combinations)

            # Generate screenshots for each combination
            for i, combination in enumerate(combinations):
                # Create new slide
                slide = prs.slides.add_slide(prs.slide_layouts[5])

                # Add title with filter values
                title = slide.shapes.title
                title_text = " | ".join([f"{f['name']}: {v}"
                                         for f, v in zip(self.filter_configs, combination)])
                title.text = title_text

                # Get and add screenshot
                image = self.get_view_image(combination)

                # Save image temporarily
                temp_image_path = f"temp_screenshot_{i}.png"
                image.save(temp_image_path)

                # Add to slide
                left = Inches(1)
                top = Inches(1.5)
                slide.shapes.add_picture(temp_image_path, left, top, height=Inches(5))

                # Clean up temp file
                os.remove(temp_image_path)

                # Update progress
                if self.progress_signal:
                    progress = int((i + 1) / total_combinations * 100)
                    self.progress_signal.emit(progress)

            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"tableau_screenshots_{timestamp}.pptx"
            prs.save(filename)

            return filename

        finally:
            if self.server:
                self.server.auth.sign_out()


def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())


if __name__ == "__main__":
    main()