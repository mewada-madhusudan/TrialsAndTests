import logging
import re
import sys
import os
import tempfile
import urllib
from datetime import datetime
from io import BytesIO
from typing import List, Dict, Tuple, Optional
import itertools
from urllib.parse import urlparse, parse_qs

import urllib3
from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QFrame, QScrollArea,
    QGroupBox, QGridLayout, QMessageBox, QProgressDialog, QComboBox
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
import tableauserverclient as TSC
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class FilterRow(QFrame):
    deleted = pyqtSignal(object)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setup_ui()

    def setup_ui(self):
        layout = QHBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        # Filter selection dropdown (will be populated later)
        self.filter_select = QComboBox()
        self.filter_select.setMinimumWidth(200)
        self.filter_select.currentIndexChanged.connect(self.on_filter_selected)

        # Options input
        self.options_input = QLineEdit()
        self.options_input.setPlaceholderText("Options (comma-separated)")
        self.options_input.setMinimumWidth(200)

        # Delete button
        delete_btn = QPushButton("×")
        delete_btn.setFixedSize(30, 30)
        delete_btn.clicked.connect(lambda: self.deleted.emit(self))

        layout.addWidget(QLabel("Filter:"))
        layout.addWidget(self.filter_select)
        layout.addWidget(QLabel("Values:"))
        layout.addWidget(self.options_input)
        layout.addWidget(delete_btn)

        self.available_values = []

    def set_available_filters(self, filters):
        """Update the filter dropdown with available filters"""
        self.filter_select.clear()
        self.filter_select.addItems([f['name'] for f in filters])

    def on_filter_selected(self, index):
        """When a filter is selected, update the options input with available values"""
        if index >= 0 and self.available_values:
            self.options_input.setText(','.join(self.available_values[index]))

    def get_data(self) -> Dict:
        return {
            'name': self.filter_select.currentText(),
            'options': [opt.strip() for opt in self.options_input.text().split(',') if opt.strip()]
        }


# class TableauViewHandler:
#     def __init__(self, server_url: str, username: str, password: str, site_id: str = ""):
#         self.server_url = self._normalize_server_url(server_url)
#         self.username = username
#         self.password = password
#         self.site_id = site_id
#         self.server = None
#
#     def _normalize_server_url(self, url: str) -> str:
#         """Ensure the server URL is properly formatted."""
#         if not url:
#             raise ValueError("Server URL cannot be empty")
#
#         if not url.startswith(('http://', 'https://')):
#             url = 'https://' + url
#
#         return url.rstrip('/')
#
#     def _extract_view_info(self, view_url: str) -> Tuple[str, str, str]:
#         """
#         Extract workbook name, view name, and/or view ID from various Tableau URL formats.
#         Returns: (workbook_name, view_name, view_id)
#         """
#         try:
#             parsed_url = urlparse(view_url)
#             path_parts = [p for p in view_url.split('/') if p]
#             query_params = parse_qs(parsed_url.query)
#
#             # Initialize return values
#             workbook_name = ""
#             view_name = ""
#             view_id = ""
#
#             # Case 1: Modern format with :view_id in the URL
#             view_id_match = re.search(r':([a-zA-Z0-9\-_]+)\?', view_url)
#             if view_id_match:
#                 view_id = view_id_match.group(1)
#                 return workbook_name, view_name, view_id
#
#             # Case 2: Older format with /views/workbook/view
#             if 'views' in path_parts:
#                 views_index = path_parts.index('views')
#                 if len(path_parts) > views_index + 2:
#                     workbook_name = path_parts[views_index + 1]
#                     view_name = path_parts[views_index + 2].split('?')[0]
#                     return workbook_name, view_name, view_id
#
#             # Case 3: Direct view ID in the last part of the URL
#             if path_parts:
#                 last_part = path_parts[-1].split('?')[0]
#                 if re.match(r'^[a-zA-Z0-9\-_]+$', last_part):
#                     view_id = last_part
#                     return workbook_name, view_name, view_id
#
#             raise ValueError("Could not extract view information from URL")
#
#         except Exception as e:
#             raise ValueError(f"Failed to parse view URL: {str(e)}")
#
#     def get_view(self, view_url: str) -> TSC.ViewItem:
#         """
#         Get the view using multiple fallback methods.
#         """
#         workbook_name, view_name, view_id = self._extract_view_info(view_url)
#
#         # Setup authentication
#         tableau_auth = TSC.TableauAuth(
#             username=self.username,
#             password=self.password,
#             site_id=self.site_id
#         )
#
#         self.server = TSC.Server(self.server_url)
#
#         try:
#             self.server.auth.sign_in(tableau_auth)
#
#             # Method 1: Try direct view ID if available
#             if view_id:
#                 try:
#                     return self.server.views.get_by_id(view_id)
#                 except:
#                     pass
#
#             # Method 2: Try getting view by name
#             if workbook_name and view_name:
#                 # Get all views
#                 all_views, pagination_item = self.server.views.get()
#
#                 # Find matching view
#                 for view in all_views:
#                     if (view.name.lower() == view_name.lower() and
#                             view.workbook.name.lower() == workbook_name.lower()):
#                         return view
#
#                 # If we need to paginate
#                 while pagination_item.has_more:
#                     all_views, pagination_item = self.server.views.get(
#                         req_options=TSC.RequestOptions(pagesize=100, pagenumber=pagination_item.page_number + 1)
#                     )
#                     for view in all_views:
#                         if (view.name.lower() == view_name.lower() and
#                                 view.workbook.name.lower() == workbook_name.lower()):
#                             return view
#
#             # Method 3: Last resort - try to match URL pattern
#             all_views, pagination_item = self.server.views.get()
#             for view in all_views:
#                 if view_url.lower().endswith(view.id.lower()):
#                     return view
#
#             raise ValueError("View not found. Please verify the URL and permissions.")
#
#         finally:
#             if self.server:
#                 self.server.auth.sign_out()
#
#     def get_available_filters(self, view: TSC.ViewItem) -> list:
#         """Get all available filters for a view."""
#         try:
#             filters = []
#             for filter_item in self.server.views.get_filters(view).all():
#                 filter_values = [item.value for item in filter_item.values]
#                 filters.append({
#                     'name': filter_item.name,
#                     'field': filter_item.field.name,
#                     'values': filter_values
#                 })
#             return filters
#         except Exception as e:
#             raise Exception(f"Failed to get filters: {str(e)}")
class TableauViewHandler:
    def __init__(self, server_url: str, username: str, password: str, site_id: str = ""):
        self.logger = None
        self.server_url = self._normalize_server_url(server_url)
        self.username = username
        self.password = password
        self.site_id = site_id
        self.server = None

    def _normalize_server_url(self, url: str) -> str:
        """Ensure the server URL is properly formatted."""
        if not url:
            raise ValueError("Server URL cannot be empty")

        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url

        return url.rstrip('/')

    def _extract_view_info(self, view_url: str) -> Tuple[str, str, str]:
        """
        Extract workbook name, view name, and/or view ID from various Tableau URL formats.
        Returns: (workbook_name, view_name, view_id)
        """
        try:
            parsed_url = urlparse(view_url)
            path_parts = [p for p in view_url.split('/') if p]
            query_params = parse_qs(parsed_url.query)

            # Initialize return values
            workbook_name = ""
            view_name = ""
            view_id = ""

            # Case 1: Modern format with :view_id in the URL
            view_id_match = re.search(r':([a-zA-Z0-9\-_]+)\?', view_url)
            if view_id_match:
                view_id = view_id_match.group(1)
                return workbook_name, view_name, view_id

            # Case 2: Older format with /views/workbook/view
            if 'views' in path_parts:
                views_index = path_parts.index('views')
                if len(path_parts) > views_index + 2:
                    workbook_name = path_parts[views_index + 1]
                    view_name = path_parts[views_index + 2].split('?')[0]
                    return workbook_name, view_name, view_id

            # Case 3: Direct view ID in the last part of the URL
            if path_parts:
                last_part = path_parts[-1].split('?')[0]
                if re.match(r'^[a-zA-Z0-9\-_]+$', last_part):
                    view_id = last_part
                    return workbook_name, view_name, view_id

            raise ValueError("Could not extract view information from URL")

        except Exception as e:
            raise ValueError(f"Failed to parse view URL: {str(e)}")

    def get_view(self, view_url: str) -> TSC.ViewItem:
        """
        Get the view using multiple fallback methods.
        """
        workbook_name, view_name, view_id = self._extract_view_info(view_url)

        # Setup authentication
        tableau_auth = TSC.TableauAuth(
            username=self.username,
            password=self.password,
            site_id=self.site_id
        )

        self.server = TSC.Server(self.server_url)

        try:
            self.server.auth.sign_in(tableau_auth)

            # Method 1: Try direct view ID if available
            if view_id:
                try:
                    return self.server.views.get_by_id(view_id)
                except:
                    pass

            # Method 2: Try getting view by name
            if workbook_name and view_name:
                # Get all views with pagination
                req_options = TSC.RequestOptions(pagesize=100)

                while True:
                    all_views, pagination_item = self.server.views.get(req_options)

                    # Find matching view
                    for view in all_views:
                        if (view.name.lower() == view_name.lower() and
                                view.workbook.name.lower() == workbook_name.lower()):
                            return view

                    # Check if there are more pages
                    if not pagination_item.total_available > (req_options.pagenumber * req_options.pagesize):
                        break

                    # Move to next page
                    req_options.pagenumber += 1

            # Method 3: Last resort - try to match URL pattern
            req_options = TSC.RequestOptions(pagesize=100)
            all_views, pagination_item = self.server.views.get(req_options)

            while True:
                for view in all_views:
                    if view_url.lower().endswith(view.id.lower()):
                        return view

                # Check if there are more pages
                if not pagination_item.total_available > (req_options.pagenumber * req_options.pagesize):
                    break

                # Move to next page
                req_options.pagenumber += 1
                all_views, pagination_item = self.server.views.get(req_options)

        except Exception as e:
            if self.server:
                self.server.auth.sign_out()
            raise ValueError("View not found. Please verify the URL and permissions.")

    def get_available_filters(self, view: TSC.ViewItem) -> list:
        """
        Get available filters for a view using methods compatible with API version 2.4.
        This version uses workbook download and metadata analysis to extract filter information.
        """
        try:
            filters = []

            # Get the workbook that contains this view
            workbook = self.server.workbooks.get_by_id(view.workbook_id)

            # Populate the workbook's views
            self.server.workbooks.populate_views(workbook)

            # Try to get workbook preview image
            try:
                self.server.workbooks.populate_preview_image(workbook)
            except Exception as e:
                print(f"Warning: Could not get workbook preview: {e}")

            # Try to get workbook permissions to ensure we have access
            try:
                self.server.workbooks.populate_permissions(workbook)
            except Exception as e:
                print(f"Warning: Could not get workbook permissions: {e}")

            # Look through the views to find matching view and extract available metadata
            for workbook_view in workbook.views:
                if workbook_view.id == view.id:
                    # Try to extract any filter information from view content url
                    if hasattr(workbook_view, 'content_url'):
                        content_url = workbook_view.content_url
                        # Parse any filter parameters from the URL
                        try:
                            parsed_url = urllib.parse.urlparse(content_url)
                            params = urllib.parse.parse_qs(parsed_url.query)
                            for param_name, param_values in params.items():
                                if param_name.startswith('vf_') or ':filter' in param_name:
                                    filter_name = param_name.replace('vf_', '').replace(':filter', '')
                                    filters.append({
                                        'name': filter_name,
                                        'field': filter_name,
                                        'values': param_values
                                    })
                        except Exception as e:
                            print(f"Warning: Could not parse content URL: {e}")

                    # Try to get any available view metadata
                    if hasattr(workbook_view, 'filter_attributes'):
                        for filter_attr in workbook_view.filter_attributes:
                            filters.append({
                                'name': filter_attr.name,
                                'field': filter_attr.name,
                                'values': filter_attr.values if hasattr(filter_attr, 'values') else []
                            })

                    break

            # If we still don't have filters, try downloading and parsing the workbook
            if not filters:
                try:
                    print("Attempting to download workbook to extract filter information...")
                    import zipfile
                    import xml.etree.ElementTree as ET

                    # Get temporary file path for workbook
                    temp_dir = tempfile.gettempdir()
                    temp_path = os.path.join(temp_dir, f"temp_workbook_{workbook.id}")

                    # Download the workbook
                    self.server.workbooks.download(workbook.id, filepath=temp_path)
                    temp_path = temp_path + ".twbx"

                    # Process the .twbx file as a zip archive
                    if os.path.exists(temp_path):
                        try:
                            with zipfile.ZipFile(temp_path, 'r') as zip_ref:
                                # Find the .twb file (XML workbook definition)
                                twb_files = [f for f in zip_ref.namelist() if f.endswith('.twb')]
                                if twb_files:
                                    # Read the first .twb file found
                                    with zip_ref.open(twb_files[0]) as twb_file:
                                        # Parse the XML content
                                        tree = ET.parse(twb_file)
                                        root = tree.getroot()

                                        # Find all filter elements
                                        for filter_elem in root.findall(".//filter"):
                                            field_name = filter_elem.get('field')
                                            if field_name:
                                                # Extract member values if available
                                                values = []
                                                for member in filter_elem.findall('.//member'):
                                                    value = member.get('value')
                                                    if value:
                                                        values.append(value)

                                                filters.append({
                                                    'name': field_name,
                                                    'field': field_name,
                                                    'values': values
                                                })

                        except Exception as e:
                            print(f"Warning: Could not parse workbook file: {e}")
                        finally:
                            # Clean up temporary file
                            try:
                                os.remove(temp_path)
                            except:
                                pass

                except Exception as e:
                    print(f"Warning: Could not download workbook: {e}")

            # Deduplicate filters based on name
            unique_filters = []
            filter_names = set()
            for f in filters:
                if f['name'] not in filter_names:
                    unique_filters.append(f)
                    filter_names.add(f['name'])

            if not unique_filters:
                print("Warning: No filters found for this view using available API methods")
                print("You may need to manually specify the filter names and values")
            else:
                print(f"Successfully found {len(unique_filters)} filters")

            return unique_filters

        except Exception as e:
            error_msg = f"Failed to get filters: {str(e)}"
            print(error_msg)
            raise Exception(error_msg)

    # def get_available_filters(self, view: TSC.ViewItem) -> list:
    #     """
    #     Get available filters for a view using methods compatible with API version 2.4.
    #     This version uses workbook download and metadata analysis to extract filter information.
    #     """
    #     try:
    #         filters = []
    #
    #         # Get the workbook that contains this view
    #         workbook = self.server.workbooks.get_by_id(view.workbook_id)
    #
    #         # Populate the workbook's views
    #         self.server.workbooks.populate_views(workbook)
    #
    #         # Try to get workbook preview image
    #         try:
    #             self.server.workbooks.populate_preview_image(workbook)
    #         except Exception as e:
    #             print(f"Warning: Could not get workbook preview: {e}")
    #
    #         # Try to get workbook permissions to ensure we have access
    #         try:
    #             self.server.workbooks.populate_permissions(workbook)
    #         except Exception as e:
    #             print(f"Warning: Could not get workbook permissions: {e}")
    #
    #         # Look through the views to find matching view and extract available metadata
    #         for workbook_view in workbook.views:
    #             if workbook_view.id == view.id:
    #                 # Try to extract any filter information from view content url
    #                 if hasattr(workbook_view, 'content_url'):
    #                     content_url = workbook_view.content_url
    #                     # Parse any filter parameters from the URL
    #                     try:
    #                         parsed_url = urllib.parse.urlparse(content_url)
    #                         params = urllib.parse.parse_qs(parsed_url.query)
    #                         for param_name, param_values in params.items():
    #                             if param_name.startswith('vf_') or ':filter' in param_name:
    #                                 filter_name = param_name.replace('vf_', '').replace(':filter', '')
    #                                 filters.append({
    #                                     'name': filter_name,
    #                                     'field': filter_name,
    #                                     'values': param_values
    #                                 })
    #                     except Exception as e:
    #                         print(f"Warning: Could not parse content URL: {e}")
    #
    #                 # Try to get any available view metadata
    #                 if hasattr(workbook_view, 'filter_attributes'):
    #                     for filter_attr in workbook_view.filter_attributes:
    #                         filters.append({
    #                             'name': filter_attr.name,
    #                             'field': filter_attr.name,
    #                             'values': filter_attr.values if hasattr(filter_attr, 'values') else []
    #                         })
    #
    #                 break
    #
    #         # If we still don't have filters, try downloading the workbook
    #         # Note: This requires appropriate permissions
    #         if not filters:
    #             try:
    #                 print("Attempting to download workbook to extract filter information...")
    #                 # Get temporary file path for workbook
    #                 temp_dir = tempfile.gettempdir()
    #                 temp_path = os.path.join(temp_dir, f"temp_workbook_{workbook.id}")
    #
    #                 # Download the workbook
    #                 self.server.workbooks.download(workbook.id, filepath=temp_path)
    #
    #                 temp_path = temp_path + ".twbx"
    #
    #                 # Parse the workbook XML to find filters
    #                 if os.path.exists(temp_path):
    #                     try:
    #                         with open(temp_path, 'r', encoding='utf-8') as f:
    #                             workbook_content = f.read()
    #
    #                         # Look for filter tags in the workbook XML
    #                         filter_matches = re.finditer(
    #                             r'<filter.*?field=[\'"](.*?)[\'"].*?</filter>',
    #                             workbook_content,
    #                             re.DOTALL
    #                         )
    #
    #                         for match in filter_matches:
    #                             field_name = match.group(1)
    #                             # Extract any available values
    #                             values = re.findall(
    #                                 r'<member value=[\'\"](.*?)[\'\"]',
    #                                 match.group(0)
    #                             )
    #
    #                             filters.append({
    #                                 'name': field_name,
    #                                 'field': field_name,
    #                                 'values': values
    #                             })
    #
    #                     except Exception as e:
    #                         print(f"Warning: Could not parse workbook file: {e}")
    #                     finally:
    #                         # Clean up temporary file
    #                         try:
    #                             os.remove(temp_path)
    #                         except:
    #                             pass
    #
    #             except Exception as e:
    #                 print(f"Warning: Could not download workbook: {e}")
    #
    #         # Deduplicate filters based on name
    #         unique_filters = []
    #         filter_names = set()
    #         for f in filters:
    #             if f['name'] not in filter_names:
    #                 unique_filters.append(f)
    #                 filter_names.add(f['name'])
    #
    #         if not unique_filters:
    #             print("Warning: No filters found for this view using available API methods")
    #             print("You may need to manually specify the filter names and values")
    #         else:
    #             print(f"Successfully found {len(unique_filters)} filters")
    #
    #         return unique_filters
    #
    #     except Exception as e:
    #         error_msg = f"Failed to get filters: {str(e)}"
    #         print(error_msg)
    #         raise Exception(error_msg)

class TableauAuthHandler:
    def __init__(self, server_url: str, username: str, password: str, site_id: str = ""):
        self.server_url = self._normalize_server_url(server_url)
        self.username = username
        self.password = password
        self.site_id = site_id
        self.server: Optional[TSC.Server] = None

        # Configure logging
        logging.basicConfig(level=logging.DEBUG)
        self.logger = logging.getLogger(__name__)

        # Disable insecure HTTPS warnings if needed
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    def _normalize_server_url(self, url: str) -> str:
        """Ensure the server URL is properly formatted."""
        if not url:
            raise ValueError("Server URL cannot be empty")

        # Add https:// if no protocol specified
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url

        # Remove trailing slashes
        url = url.rstrip('/')

        # Validate URL format
        try:
            parsed = urlparse(url)
            if not parsed.netloc:
                raise ValueError("Invalid server URL format")
        except Exception as e:
            raise ValueError(f"Invalid server URL: {str(e)}")

        return url

    def authenticate(self) -> Tuple[TSC.Server, str]:
        """
        Attempt to authenticate with Tableau Server.
        Returns tuple of (server, auth_token)
        """
        try:
            # Initialize auth object with retry settings
            tableau_auth = TSC.TableauAuth(
                username=self.username,
                password=self.password,
                site_id=self.site_id
            )

            # Configure server with additional options
            self.server = TSC.Server(
                self.server_url,
                use_server_version=True,
                http_options={
                    'verify': False,  # Consider making this configurable
                    'timeout': 30
                }
            )

            # Attempt authentication
            self.logger.info(f"Attempting to authenticate to {self.server_url}")
            auth_token = self.server.auth.sign_in(tableau_auth)

            self.logger.info("Authentication successful")
            return self.server, auth_token

        except TSC.ServerResponseError as e:
            self.logger.error(f"Server response error: {str(e)}")
            if hasattr(e, 'response') and e.response:
                self.logger.error(f"Response content: {e.response.content}")
            raise

        except Exception as e:
            self.logger.error(f"Authentication failed: {str(e)}")
            raise

    def sign_out(self):
        """Safely sign out from the server."""
        try:
            if self.server:
                self.server.auth.sign_out()
                self.logger.info("Successfully signed out")
        except Exception as e:
            self.logger.error(f"Error during sign out: {str(e)}")

    def __enter__(self):
        """Enable context manager usage."""
        server, token = self.authenticate()
        return server, token

    def __exit__(self, exc_type, exc_val, exc_tb):
        """Ensure proper cleanup when used as context manager."""
        self.sign_out()

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.filter_rows: List[FilterRow] = []
        self.available_filters = []
        self.setup_ui()

    def setup_ui(self):
        self.setWindowTitle("Tableau Dashboard Screenshot Generator")
        self.setMinimumSize(800, 600)

        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)

        # Authentication group
        auth_group = QGroupBox("Tableau Authentication")
        auth_layout = QGridLayout()

        self.server_url = QLineEdit()
        self.server_url.setPlaceholderText("https://your-tableau-server")
        self.username = QLineEdit()
        self.username.setPlaceholderText("username")
        self.password = QLineEdit()
        self.password.setPlaceholderText("password")
        self.password.setEchoMode(QLineEdit.EchoMode.Password)
        self.site_id = QLineEdit()
        self.site_id.setPlaceholderText("site-id (leave empty for default)")

        auth_layout.addWidget(QLabel("Server URL:"), 0, 0)
        auth_layout.addWidget(self.server_url, 0, 1)
        auth_layout.addWidget(QLabel("Username:"), 1, 0)
        auth_layout.addWidget(self.username, 1, 1)
        auth_layout.addWidget(QLabel("Password:"), 2, 0)
        auth_layout.addWidget(self.password, 2, 1)
        auth_layout.addWidget(QLabel("Site ID:"), 3, 0)
        auth_layout.addWidget(self.site_id, 3, 1)

        auth_group.setLayout(auth_layout)
        main_layout.addWidget(auth_group)

        # Dashboard URL group
        url_group = QGroupBox("Dashboard View URL")
        url_layout = QVBoxLayout()
        self.view_url = QLineEdit()
        self.view_url.setPlaceholderText("Enter the Tableau dashboard view URL")

        # Add Load Filters button
        load_filters_btn = QPushButton("Load Available Filters")
        load_filters_btn.clicked.connect(self.load_available_filters)

        url_layout.addWidget(self.view_url)
        url_layout.addWidget(load_filters_btn)
        url_group.setLayout(url_layout)
        main_layout.addWidget(url_group)

        # Filters group
        filters_group = QGroupBox("Filters")
        filters_layout = QVBoxLayout()

        # Scroll area for filters
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll_widget = QWidget()
        self.filters_layout = QVBoxLayout(scroll_widget)
        scroll.setWidget(scroll_widget)

        add_filter_btn = QPushButton("Add Filter")
        add_filter_btn.clicked.connect(self.add_filter_row)
        filters_layout.addWidget(add_filter_btn)
        filters_layout.addWidget(scroll)

        filters_group.setLayout(filters_layout)
        main_layout.addWidget(filters_group)

        # Generate button
        generate_btn = QPushButton("Generate PPT")
        generate_btn.setStyleSheet("background-color: #4CAF50; color: white; padding: 10px;")
        generate_btn.clicked.connect(self.generate_screenshots)
        main_layout.addWidget(generate_btn)

    # Modified MainWindow authentication method
    def authenticate_tableau(self):
        """
        Modified authentication method for the MainWindow class.
        """
        auth_config = self.get_auth_config()

        try:
            auth_handler = TableauAuthHandler(
                server_url=auth_config['server_url'],
                username=auth_config['username'],
                password=auth_config['password'],
                site_id=auth_config['site_id']
            )

            # Use context manager for automatic cleanup
            with auth_handler as (server, auth_token):
                # Store server instance for later use
                self.tableau_server = server
                return True

        except ValueError as e:
            QMessageBox.critical(self, "Configuration Error", str(e))
            return False

        except TSC.ServerResponseError as e:
            QMessageBox.critical(self, "Server Error",
                                 f"Failed to connect to Tableau server: {str(e)}\n"
                                 f"Please verify your server URL and credentials.")
            return False

        except Exception as e:
            QMessageBox.critical(self, "Authentication Error",
                                 f"An unexpected error occurred: {str(e)}")
            return False

    def load_available_filters(self):
        """Load available filters from the dashboard"""
        if not all([self.server_url.text(), self.username.text(), self.password.text()]):
            QMessageBox.warning(self, "Error", "Please fill in all authentication fields.")
            return

        if not self.view_url.text():
            QMessageBox.warning(self, "Error", "Please enter a dashboard view URL.")
            return

        try:
            view_handler = TableauViewHandler(
                server_url=self.server_url.text(),
                username=self.username.text(),
                password=self.password.text(),
                site_id=self.site_id.text()
            )

            # Get the view using the enhanced handler
            view = view_handler.get_view(self.view_url.text())

            # Get and store available filters
            self.available_filters = view_handler.get_available_filters(view)

            # Update existing filter rows
            for row in self.filter_rows:
                row.set_available_filters(self.available_filters)
                row.available_values = [f['values'] for f in self.available_filters]

            QMessageBox.information(self, "Success", f"Loaded {len(self.available_filters)} filters!")

        except ValueError as e:
            QMessageBox.critical(self, "Error", str(e))
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to load filters: {str(e)}")

    def add_filter_row(self):
        if not self.available_filters:
            QMessageBox.warning(self, "Error", "Please load available filters first.")
            return

        filter_row = FilterRow()
        filter_row.set_available_filters(self.available_filters)
        filter_row.available_values = [f['values'] for f in self.available_filters]
        filter_row.deleted.connect(self.remove_filter_row)
        self.filters_layout.addWidget(filter_row)
        self.filter_rows.append(filter_row)

    def remove_filter_row(self, row):
        self.filters_layout.removeWidget(row)
        self.filter_rows.remove(row)
        row.deleteLater()

    def get_auth_config(self):
        return {
            'server_url': self.server_url.text(),
            'username': self.username.text(),
            'password': self.password.text(),
            'site_id': self.site_id.text()
        }

    def generate_screenshots(self):
        if not self.filter_rows:
            QMessageBox.warning(self, "Error", "Please add at least one filter.")
            return

        # Get filter configurations
        filters = []
        for row in self.filter_rows:
            filter_data = row.get_data()
            filter_index = row.filter_select.currentIndex()
            if filter_index >= 0:
                filter_data['field'] = self.available_filters[filter_index]['field']
                filters.append(filter_data)

        # Start generation in a separate thread
        self.generator_thread = GeneratorThread(
            self.get_auth_config(),
            self.view_url.text(),
            filters
        )

        # Show progress dialog
        self.progress_dialog = QProgressDialog("Generating screenshots...", "Cancel", 0, 100, self)
        self.progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)

        # Connect signals
        self.generator_thread.progress.connect(self.progress_dialog.setValue)
        self.generator_thread.finished.connect(self.generation_finished)
        self.generator_thread.error.connect(self.generation_error)

        self.generator_thread.start()

    def generation_finished(self, filename):
        self.progress_dialog.close()
        QMessageBox.information(self, "Success",
                                f"PowerPoint has been generated successfully!\nSaved as: {filename}")

    def generation_error(self, error_message):
        self.progress_dialog.close()
        QMessageBox.critical(self, "Error", f"An error occurred: {error_message}")


class TableauScreenshotGenerator:
    def __init__(self, auth_config, view_url, filter_configs):
        self.auth_config = auth_config
        self.view_url = view_url
        self.filter_configs = filter_configs
        self.progress_signal = None
        self.server = None

    def setup_tableau_client(self):
        self.tableau_auth = TSC.TableauAuth(
            username=self.auth_config['username'],
            password=self.auth_config['password'],
            site_id=self.auth_config['site_id']
        )
        self.server = TSC.Server(self.auth_config['server_url'])

    def get_view_image(self, view, filter_values):
        """Get dashboard screenshot with specified filter values"""
        try:
            # Apply filters
            for filter_config, value in filter_values:
                self.server.views.update_filter(
                    view,
                    filter_config['field'],
                    [value]
                )

            # Get image after filters are applied
            image_response = self.server.views.get_image(view)
            return Image.open(BytesIO(image_response))
        except Exception as e:
            raise Exception(f"Failed to get screenshot: {str(e)}")

    def generate_presentation(self):
        """Generate PowerPoint with screenshots for all filter combinations"""
        try:
            self.setup_tableau_client()
            self.server.auth.sign_in(self.tableau_auth)

            # Get view
            view_id = self.view_url.split('/')[-1]
            view = self.server.views.get_by_id(view_id)

            prs = Presentation()

            # Calculate all possible filter combinations
            filter_options = [config['options'] for config in self.filter_configs]
            combinations = list(itertools.product(*filter_options))
            total_combinations = len(combinations)

            # Generate screenshot for each combination
            for i, combination in enumerate(combinations):
                # Create new slide
                slide = prs.slides.add_slide(prs.slide_layouts[5])

                # Add title showing filter values
                title = slide.shapes.title
                title_text = " | ".join(
                    f"{f['name']}: {v}"
                    for f, v in zip(self.filter_configs, combination)
                )
                title.text = title_text

                # Get and add screenshot
                filter_values = list(zip(self.filter_configs, combination))
                image = self.get_view_image(view, filter_values)

                # Save image temporarily
                temp_path = f"temp_screenshot_{i}.png"
                image.save(temp_path)

                # Add to slide and clean up
                slide.shapes.add_picture(
                    temp_path,
                    left=Inches(1),
                    top=Inches(1.5),
                    height=Inches(5)
                )
                os.remove(temp_path)

                # Update progress
                if self.progress_signal:
                    progress = int((i + 1) / total_combinations * 100)
                    self.progress_signal.emit(progress)

            # Save presentation with timestamp
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"tableau_screenshots_{timestamp}.pptx"
            prs.save(filename)
            return filename

        finally:
            if self.server:
                self.server.auth.sign_out()


class GeneratorThread(QThread):
    progress = pyqtSignal(int)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, auth_config, view_url, filters):
        super().__init__()
        self.auth_config = auth_config
        self.view_url = view_url
        self.filters = filters

    def run(self):
        try:
            generator = TableauScreenshotGenerator(
                self.auth_config,
                self.view_url,
                self.filters
            )
            generator.progress_signal = self.progress
            filename = generator.generate_presentation()
            self.finished.emit(filename)
        except Exception as e:
            self.error.emit(str(e))


def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())


if __name__ == "__main__":
    main()
# https://prod-apnortheast-a.online.tableau.com/t/madhusudan145-67ed7c4919/view/sample_dash/Sheet1?iid=76c71b72-02bb-49d2-a7e8-ba74d1a0deca