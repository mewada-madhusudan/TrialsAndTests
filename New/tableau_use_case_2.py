import sys
import os
from datetime import datetime
from io import BytesIO
from typing import List, Dict
import itertools

from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QFrame, QScrollArea,
    QGroupBox, QGridLayout, QMessageBox, QProgressDialog, QComboBox
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
import tableauserverclient as TSC
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class FilterRow(QFrame):
    deleted = pyqtSignal(object)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setup_ui()

    def setup_ui(self):
        layout = QHBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        # Filter selection dropdown (will be populated later)
        self.filter_select = QComboBox()
        self.filter_select.setMinimumWidth(200)
        self.filter_select.currentIndexChanged.connect(self.on_filter_selected)

        # Options input
        self.options_input = QLineEdit()
        self.options_input.setPlaceholderText("Options (comma-separated)")
        self.options_input.setMinimumWidth(200)

        # Delete button
        delete_btn = QPushButton("×")
        delete_btn.setFixedSize(30, 30)
        delete_btn.clicked.connect(lambda: self.deleted.emit(self))

        layout.addWidget(QLabel("Filter:"))
        layout.addWidget(self.filter_select)
        layout.addWidget(QLabel("Values:"))
        layout.addWidget(self.options_input)
        layout.addWidget(delete_btn)

        self.available_values = []

    def set_available_filters(self, filters):
        """Update the filter dropdown with available filters"""
        self.filter_select.clear()
        self.filter_select.addItems([f['name'] for f in filters])

    def on_filter_selected(self, index):
        """When a filter is selected, update the options input with available values"""
        if index >= 0 and self.available_values:
            self.options_input.setText(','.join(self.available_values[index]))

    def get_data(self) -> Dict:
        return {
            'name': self.filter_select.currentText(),
            'options': [opt.strip() for opt in self.options_input.text().split(',') if opt.strip()]
        }


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.filter_rows: List[FilterRow] = []
        self.available_filters = []
        self.setup_ui()

    def setup_ui(self):
        self.setWindowTitle("Tableau Dashboard Screenshot Generator")
        self.setMinimumSize(800, 600)

        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)

        # Authentication group
        auth_group = QGroupBox("Tableau Authentication")
        auth_layout = QGridLayout()

        self.server_url = QLineEdit()
        self.server_url.setPlaceholderText("https://your-tableau-server")
        self.username = QLineEdit()
        self.username.setPlaceholderText("username")
        self.password = QLineEdit()
        self.password.setPlaceholderText("password")
        self.password.setEchoMode(QLineEdit.EchoMode.Password)
        self.site_id = QLineEdit()
        self.site_id.setPlaceholderText("site-id (leave empty for default)")

        auth_layout.addWidget(QLabel("Server URL:"), 0, 0)
        auth_layout.addWidget(self.server_url, 0, 1)
        auth_layout.addWidget(QLabel("Username:"), 1, 0)
        auth_layout.addWidget(self.username, 1, 1)
        auth_layout.addWidget(QLabel("Password:"), 2, 0)
        auth_layout.addWidget(self.password, 2, 1)
        auth_layout.addWidget(QLabel("Site ID:"), 3, 0)
        auth_layout.addWidget(self.site_id, 3, 1)

        auth_group.setLayout(auth_layout)
        main_layout.addWidget(auth_group)

        # Dashboard URL group
        url_group = QGroupBox("Dashboard View URL")
        url_layout = QVBoxLayout()
        self.view_url = QLineEdit()
        self.view_url.setPlaceholderText("Enter the Tableau dashboard view URL")

        # Add Load Filters button
        load_filters_btn = QPushButton("Load Available Filters")
        load_filters_btn.clicked.connect(self.load_available_filters)

        url_layout.addWidget(self.view_url)
        url_layout.addWidget(load_filters_btn)
        url_group.setLayout(url_layout)
        main_layout.addWidget(url_group)

        # Filters group
        filters_group = QGroupBox("Filters")
        filters_layout = QVBoxLayout()

        # Scroll area for filters
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll_widget = QWidget()
        self.filters_layout = QVBoxLayout(scroll_widget)
        scroll.setWidget(scroll_widget)

        add_filter_btn = QPushButton("Add Filter")
        add_filter_btn.clicked.connect(self.add_filter_row)
        filters_layout.addWidget(add_filter_btn)
        filters_layout.addWidget(scroll)

        filters_group.setLayout(filters_layout)
        main_layout.addWidget(filters_group)

        # Generate button
        generate_btn = QPushButton("Generate PPT")
        generate_btn.setStyleSheet("background-color: #4CAF50; color: white; padding: 10px;")
        generate_btn.clicked.connect(self.generate_screenshots)
        main_layout.addWidget(generate_btn)

    def load_available_filters(self):
        """Load available filters from the dashboard"""
        if not all([self.server_url.text(), self.username.text(), self.password.text()]):
            QMessageBox.warning(self, "Error", "Please fill in all authentication fields.")
            return

        if not self.view_url.text():
            QMessageBox.warning(self, "Error", "Please enter a dashboard view URL.")
            return

        try:
            # Setup Tableau client
            tableau_auth = TSC.TableauAuth(
                username=self.username.text(),
                password=self.password.text(),
                site_id=self.site_id.text()
            )
            server = TSC.Server(self.server_url.text())

            # Sign in and get view
            server.auth.sign_in(tableau_auth)
            view_id = self.view_url.text().split('/')[-1]
            view = server.views.get_by_id(view_id)

            # Get available filters
            self.available_filters = []
            for filter_item in server.views.get_filters(view).all():
                filter_values = [item.value for item in filter_item.values]
                self.available_filters.append({
                    'name': filter_item.name,
                    'field': filter_item.field.name,
                    'values': filter_values
                })

            # Update existing filter rows
            for row in self.filter_rows:
                row.set_available_filters(self.available_filters)
                row.available_values = [f['values'] for f in self.available_filters]

            QMessageBox.information(self, "Success", f"Loaded {len(self.available_filters)} filters!")

        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to load filters: {str(e)}")
        finally:
            if server:
                server.auth.sign_out()

    def add_filter_row(self):
        if not self.available_filters:
            QMessageBox.warning(self, "Error", "Please load available filters first.")
            return

        filter_row = FilterRow()
        filter_row.set_available_filters(self.available_filters)
        filter_row.available_values = [f['values'] for f in self.available_filters]
        filter_row.deleted.connect(self.remove_filter_row)
        self.filters_layout.addWidget(filter_row)
        self.filter_rows.append(filter_row)

    def remove_filter_row(self, row):
        self.filters_layout.removeWidget(row)
        self.filter_rows.remove(row)
        row.deleteLater()

    def get_auth_config(self):
        return {
            'server_url': self.server_url.text(),
            'username': self.username.text(),
            'password': self.password.text(),
            'site_id': self.site_id.text()
        }

    def generate_screenshots(self):
        if not self.filter_rows:
            QMessageBox.warning(self, "Error", "Please add at least one filter.")
            return

        # Get filter configurations
        filters = []
        for row in self.filter_rows:
            filter_data = row.get_data()
            filter_index = row.filter_select.currentIndex()
            if filter_index >= 0:
                filter_data['field'] = self.available_filters[filter_index]['field']
                filters.append(filter_data)

        # Start generation in a separate thread
        self.generator_thread = GeneratorThread(
            self.get_auth_config(),
            self.view_url.text(),
            filters
        )

        # Show progress dialog
        self.progress_dialog = QProgressDialog("Generating screenshots...", "Cancel", 0, 100, self)
        self.progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)

        # Connect signals
        self.generator_thread.progress.connect(self.progress_dialog.setValue)
        self.generator_thread.finished.connect(self.generation_finished)
        self.generator_thread.error.connect(self.generation_error)

        self.generator_thread.start()

    def generation_finished(self, filename):
        self.progress_dialog.close()
        QMessageBox.information(self, "Success",
                                f"PowerPoint has been generated successfully!\nSaved as: {filename}")

    def generation_error(self, error_message):
        self.progress_dialog.close()
        QMessageBox.critical(self, "Error", f"An error occurred: {error_message}")


class TableauScreenshotGenerator:
    def __init__(self, auth_config, view_url, filter_configs):
        self.auth_config = auth_config
        self.view_url = view_url
        self.filter_configs = filter_configs
        self.progress_signal = None
        self.server = None

    def setup_tableau_client(self):
        self.tableau_auth = TSC.TableauAuth(
            username=self.auth_config['username'],
            password=self.auth_config['password'],
            site_id=self.auth_config['site_id']
        )
        self.server = TSC.Server(self.auth_config['server_url'])

    def get_view_image(self, view, filter_values):
        """Get dashboard screenshot with specified filter values"""
        try:
            # Apply filters
            for filter_config, value in filter_values:
                self.server.views.update_filter(
                    view,
                    filter_config['field'],
                    [value]
                )

            # Get image after filters are applied
            image_response = self.server.views.get_image(view)
            return Image.open(BytesIO(image_response))
        except Exception as e:
            raise Exception(f"Failed to get screenshot: {str(e)}")

    def generate_presentation(self):
        """Generate PowerPoint with screenshots for all filter combinations"""
        try:
            self.setup_tableau_client()
            self.server.auth.sign_in(self.tableau_auth)

            # Get view
            view_id = self.view_url.split('/')[-1]
            view = self.server.views.get_by_id(view_id)

            prs = Presentation()

            # Calculate all possible filter combinations
            filter_options = [config['options'] for config in self.filter_configs]
            combinations = list(itertools.product(*filter_options))
            total_combinations = len(combinations)

            # Generate screenshot for each combination
            for i, combination in enumerate(combinations):
                # Create new slide
                slide = prs.slides.add_slide(prs.slide_layouts[5])

                # Add title showing filter values
                title = slide.shapes.title
                title_text = " | ".join(
                    f"{f['name']}: {v}"
                    for f, v in zip(self.filter_configs, combination)
                )
                title.text = title_text

                # Get and add screenshot
                filter_values = list(zip(self.filter_configs, combination))
                image = self.get_view_image(view, filter_values)

                # Save image temporarily
                temp_path = f"temp_screenshot_{i}.png"
                image.save(temp_path)

                # Add to slide and clean up
                slide.shapes.add_picture(
                    temp_path,
                    left=Inches(1),
                    top=Inches(1.5),
                    height=Inches(5)
                )
                os.remove(temp_path)

                # Update progress
                if self.progress_signal:
                    progress = int((i + 1) / total_combinations * 100)
                    self.progress_signal.emit(progress)

            # Save presentation with timestamp
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"tableau_screenshots_{timestamp}.pptx"
            prs.save(filename)
            return filename

        finally:
            if self.server:
                self.server.auth.sign_out()


class GeneratorThread(QThread):
    progress = pyqtSignal(int)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, auth_config, view_url, filters):
        super().__init__()
        self.auth_config = auth_config
        self.view_url = view_url
        self.filters = filters

    def run(self):
        try:
            generator = TableauScreenshotGenerator(
                self.auth_config,
                self.view_url,
                self.filters
            )
            generator.progress_signal = self.progress
            filename = generator.generate_presentation()
            self.finished.emit(filename)
        except Exception as e:
            self.error.emit(str(e))


def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())


if __name__ == "__main__":
    main()