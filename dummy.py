import os
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import tableauserverclient as TSC
from pptx import Presentation
from pptx.util import Inches
import itertools
import logging
import io
from PIL import Image

class TableauServerProcessor:
    def __init__(self):
        """
        Initialize Tableau Server connection and logging
        """
        # Logging setup
        logging.basicConfig(level=logging.INFO, 
                            format='%(asctime)s - %(levelname)s: %(message)s')
        self.logger = logging.getLogger(__name__)
        
        # Tableau Server connection variables
        self.server = None
        self.tableau_auth = None
        self.site = None
        
        # Workbook and filter details
        self.selected_workbook = None
        self.selected_view = None
        self.filters = {}
        
        # Main application window
        self.root = None
        
        # Filter selection variables
        self.filter_vars = {}
    
    def connect(self, server_url, username, password, site_id=''):
        """
        Connect to Tableau Server
        
        :param server_url: Tableau Server URL
        :param username: Tableau username
        :param password: Tableau password
        :param site_id: Optional site ID
        """
        try:
            # Create Tableau authentication object
            self.tableau_auth = TSC.TableauAuth(username, password, site_id)
            
            # Create server connection
            self.server = TSC.Server(server_url)
            
            # Sign in to Tableau Server
            with self.server.auth.sign_in(self.tableau_auth):
                self.site = self.server.site
                self.logger.info("Successfully connected to Tableau Server")
                return True
        
        except TSC.ServerResponseError as e:
            self.logger.error(f"Authentication failed: {e}")
            messagebox.showerror("Authentication Error", str(e))
            return False
    
    def get_workbooks(self):
        """
        Retrieve list of workbooks for the authenticated user
        
        :return: List of workbooks
        """
        try:
            # Retrieve all workbooks 
            with self.server.auth.sign_in(self.tableau_auth):
                all_workbooks = list(TSC.Endpoint.get(self.server.workbooks))
                return all_workbooks
        
        except TSC.ServerResponseError as e:
            self.logger.error(f"Failed to retrieve workbooks: {e}")
            messagebox.showerror("Workbook Retrieval Error", str(e))
            return []
    
    def get_views_for_workbook(self, workbook):
        """
        Get views for a specific workbook
        
        :param workbook: Tableau workbook object
        :return: List of views
        """
        try:
            # Reload workbook to get views
            with self.server.auth.sign_in(self.tableau_auth):
                self.server.workbooks.populate_views(workbook)
                return list(workbook.views)
        
        except TSC.ServerResponseError as e:
            self.logger.error(f"Failed to retrieve views: {e}")
            messagebox.showerror("View Retrieval Error", str(e))
            return []
    
    def get_worksheet_filters(self, view):
        """
        Get filters for a specific view
        
        :param view: Tableau view object
        :return: Dictionary of filters
        """
        try:
            # Populate view's datasources to access fields
            with self.server.auth.sign_in(self.tableau_auth):
                self.server.views.populate_datasources(view)
                
                filters = {}
                for datasource in view.datasources:
                    # Populate datasource's fields
                    self.server.datasources.populate_fields(datasource)
                    
                    # Extract unique filter fields
                    for field in datasource.fields:
                        if field.is_dimension:
                            # Fetch unique values for the field
                            unique_values = self._get_unique_field_values(datasource, field)
                            if unique_values:
                                filters[field.name] = unique_values
                
                return filters
        
        except TSC.ServerResponseError as e:
            self.logger.error(f"Failed to retrieve filters: {e}")
            return {}
    
    def _get_unique_field_values(self, datasource, field, limit=50):
        """
        Get unique values for a field
        
        :param datasource: Tableau datasource
        :param field: Field to get values for
        :param limit: Maximum number of values to retrieve
        :return: List of unique values
        """
        try:
            # In tableauserverclient, getting unique values might require 
            # additional data source queries which are not directly supported
            # This is a placeholder approach
            return []  # Modify as per your specific Tableau Server setup
        
        except Exception as e:
            self.logger.error(f"Failed to get unique field values: {e}")
            return []
    
    def capture_view_image(self, view, filter_combination=None):
        """
        Capture image of the view with optional filters
        
        :param view: Tableau view object
        :param filter_combination: Optional dictionary of filter values
        :return: PIL Image object
        """
        try:
            with self.server.auth.sign_in(self.tableau_auth):
                # Download view image
                image_req = self.server.views.get_image(view)
                image = Image.open(io.BytesIO(image_req))
                return image
        
        except TSC.ServerResponseError as e:
            self.logger.error(f"Failed to capture view image: {e}")
            return None
    
    def create_ui(self):
        """
        Create user interface for workbook and view selection
        """
        self.root = tk.Tk()
        self.root.title("Tableau Dashboard Filter Processor")
        self.root.geometry("600x500")
        
        # Workbook selection
        tk.Label(self.root, text="Select Workbook:").pack(pady=(10,0))
        workbook_listbox = tk.Listbox(self.root, width=50)
        workbook_listbox.pack(pady=5)
        
        # Populate workbooks
        workbooks = self.get_workbooks()
        for wb in workbooks:
            workbook_listbox.insert(tk.END, f"{wb.name}")
        
        # View selection
        tk.Label(self.root, text="Select View:").pack(pady=(10,0))
        view_listbox = tk.Listbox(self.root, width=50)
        view_listbox.pack(pady=5)
        
        def on_workbook_select(event):
            # Clear previous view selections
            view_listbox.delete(0, tk.END)
            
            # Get selected workbook
            selection = workbook_listbox.curselection()
            if not selection:
                return
            
            selected_wb = workbooks[selection[0]]
            
            # Get views for selected workbook
            views = self.get_views_for_workbook(selected_wb)
            
            for view in views:
                view_listbox.insert(tk.END, f"{view.name}")
        
        def on_view_select(event):
            # Get selected view
            wb_selection = workbook_listbox.curselection()
            view_selection = view_listbox.curselection()
            
            if not wb_selection or not view_selection:
                return
            
            selected_wb = workbooks[wb_selection[0]]
            views = self.get_views_for_workbook(selected_wb)
            selected_view = views[view_selection[0]]
            
            self.selected_workbook = selected_wb
            self.selected_view = selected_view
            
            # Fetch filters for the selected view
            self.filters = self.get_worksheet_filters(selected_view)
            
            # Close current window and open filter selection
            self.root.destroy()
            self.create_filter_selection_ui()
        
        # Bind selection events
        workbook_listbox.bind('<<ListboxSelect>>', on_workbook_select)
        view_listbox.bind('<<ListboxSelect>>', on_view_select)
        
        self.root.mainloop()
    
    def create_filter_selection_ui(self):
        """
        Create UI for selecting filter combinations
        """
        self.root = tk.Tk()
        self.root.title("Tableau Filter Selection")
        self.root.geometry("600x500")
        
        # Filter selection frames
        for filter_name, options in self.filters.items():
            frame = ttk.LabelFrame(self.root, text=filter_name)
            frame.pack(padx=10, pady=10, fill='x')
            
            # Multiselect listbox
            self.filter_vars[filter_name] = []
            listbox = tk.Listbox(frame, selectmode=tk.MULTIPLE, height=5)
            for option in options:
                listbox.insert(tk.END, option)
            listbox.pack(padx=10, pady=10, fill='x')
            
            # Bind selection
            listbox.bind('<<ListboxSelect>>', 
                         lambda e, name=filter_name: self.update_filter_selection(e, name))
        
        # Generate button
        generate_btn = ttk.Button(self.root, text="Generate PowerPoint", command=self.generate_ppt)
        generate_btn.pack(pady=20)
        
        self.root.mainloop()
    
    def update_filter_selection(self, event, filter_name):
        """
        Update selected filter options
        
        :param event: Tkinter event
        :param filter_name: Name of the filter
        """
        listbox = event.widget
        selected_indices = listbox.curselection()
        
        # Update filter selections
        self.filter_vars[filter_name] = [
            listbox.get(idx) for idx in selected_indices
        ]
    
    def generate_ppt(self):
        """
        Generate PowerPoint with all filter combinations
        """
        # Validate filter selections
        for filter_name, selections in self.filter_vars.items():
            if not selections:
                messagebox.showerror("Error", f"Please select at least one option for {filter_name}")
                return
        
        # Generate filter combinations
        combinations = list(itertools.product(
            *[self.filter_vars[filter_name] for filter_name in self.filters.keys()]
        ))
        
        # Choose save location
        save_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        
        if not save_path:
            return
        
        # Create PowerPoint
        prs = Presentation()
        
        # Add slides for each combination
        for i, combo in enumerate(combinations, 1):
            # Create combination dictionary
            combo_dict = dict(zip(self.filters.keys(), combo))
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            
            # Add text with filter combination details
            title = slide.shapes.title
            title.text = f"Combination {i}"
            
            # Add text box with filter details
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(6), Inches(2)
            )
            tf = txBox.text_frame
            
            # Add filter details to text box
            for j, (filter_name, value) in enumerate(combo_dict.items()):
                p = tf.add_paragraph()
                p.text = f"{filter_name}: {value}"
            
            # Capture view image
            image = self.capture_view_image(self.selected_view)
            if image:
                # Add image to slide
                pic_left = Inches(1)
                pic_top = Inches(4)
                pic_width = Inches(6)
                
                # Save image temporarily
                temp_image_path = 'temp_view_image.png'
                image.save(temp_image_path)
                
                # Add picture to slide
                slide.shapes.add_picture(
                    temp_image_path, pic_left, pic_top, pic_width
                )
                
                # Clean up temporary file
                os.remove(temp_image_path)
        
        # Save PowerPoint
        prs.save(save_path)
        
        messagebox.showinfo("Success", "PowerPoint generated successfully!")

def main():
    # Create root login window
    login_window = tk.Tk()
    login_window.title("Tableau Server Login")
    login_window.geometry("400x300")

    # Server URL
    tk.Label(login_window, text="Server URL:").pack(pady=(10,0))
    server_entry = tk.Entry(login_window, width=40)
    server_entry.pack(pady=5)
    server_entry.insert(0, "https://your-tableau-server.com")

    # Username
    tk.Label(login_window, text="Username:").pack(pady=(10,0))
    username_entry = tk.Entry(login_window, width=40)
    username_entry.pack(pady=5)

    # Password
    tk.Label(login_window, text="Password:").pack(pady=(10,0))
    password_entry = tk.Entry(login_window, show="*", width=40)
    password_entry.pack(pady=5)

    # Site (optional)
    tk.Label(login_window, text="Site ID (optional):").pack(pady=(10,0))
    site_entry = tk.Entry(login_window, width=40)
    site_entry.pack(pady=5)

    def login():
        # Get login details
        server_url = server_entry.get()
        username = username_entry.get()
        password = password_entry.get()
        site_id = site_entry.get()

        # Create processor
        processor = TableauServerProcessor()

        # Attempt to connect
        if processor.connect(server_url, username, password, site_id):
            login_window.destroy()
            processor.create_ui()

    # Login button
    login_btn = tk.Button(login_window, text="Login", command=login)
    login_btn.pack(pady=20)

    login_window.mainloop()

if __name__ == "__main__":
    main()
