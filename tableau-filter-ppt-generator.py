import os
import requests
from urllib.parse import urlparse, parse_qs
from pptx import Presentation
from pptx.util import Inches

# Alternative web automation libraries
try:
    from playwright.sync_api import sync_playwright
except ImportError:
    sync_playwright = None

try:
    import pyppeteer
except ImportError:
    pyppeteer = None

class TableauFilterPPTGenerator:
    def __init__(self, tableau_view_url):
        """
        Initialize the Tableau Filter PPT Generator
        
        :param tableau_view_url: Full URL of the Tableau view
        """
        self.tableau_view_url = tableau_view_url
        self.driver = None
        self.filters = {}
        self.playwright_context = None
    
    def connect_with_playwright(self, username=None, password=None):
        """
        Connect to Tableau view using Playwright
        
        :param username: Optional Tableau username
        :param password: Optional Tableau password
        """
        if not sync_playwright:
            raise ImportError("Playwright is not installed. Please install it using 'pip install playwright'")
        
        self.playwright_context = sync_playwright().start()
        browser = self.playwright_context.chromium.launch(headless=False)
        page = browser.new_page()
        
        # Navigate to Tableau view
        page.goto(self.tableau_view_url)
        
        # Handle authentication if credentials provided
        if username and password:
            # Adjust these selectors based on your Tableau login page
            page.fill('#username', username)
            page.fill('#password', password)
            page.click('#login-button')
        
        self.driver = page
        return page
    
    def connect_with_pyppeteer(self, username=None, password=None):
        """
        Connect to Tableau view using Pyppeteer
        
        :param username: Optional Tableau username
        :param password: Optional Tableau password
        """
        if not pyppeteer:
            raise ImportError("Pyppeteer is not installed. Please install it using 'pip install pyppeteer'")
        
        async def setup_browser():
            browser = await pyppeteer.launch(headless=False)
            page = await browser.newPage()
            await page.goto(self.tableau_view_url)
            
            # Handle authentication if credentials provided
            if username and password:
                # Adjust these selectors based on your Tableau login page
                await page.type('#username', username)
                await page.type('#password', password)
                await page.click('#login-button')
            
            return page, browser
        
        import asyncio
        page, self.browser = asyncio.get_event_loop().run_until_complete(setup_browser())
        self.driver = page
        return page
    
    def extract_filters_with_requests(self):
        """
        Extract filters using requests (limited functionality)
        
        :return: Dictionary of potential filters
        """
        try:
            response = requests.get(self.tableau_view_url)
            response.raise_for_status()
            
            # This is a very basic extraction and might not work for all Tableau views
            import re
            filter_matches = re.findall(r'data-filter-name="([^"]+)"', response.text)
            
            self.filters = {
                filter_name: ['Option 1', 'Option 2']  # Placeholder
                for filter_name in set(filter_matches)
            }
            
            return self.filters
        except Exception as e:
            print(f"Error extracting filters: {e}")
            return {}
    
    def apply_filters(self, selected_filters):
        """
        Apply selected filters (method depends on connection type)
        
        :param selected_filters: Dictionary of {filter_name: [selected_options]}
        """
        if not self.driver:
            raise ValueError("No web driver connected. Use connect_with_playwright() or connect_with_pyppeteer() first.")
        
        # Playwright method
        if hasattr(self.driver, 'fill'):
            for filter_name, selected_options in selected_filters.items():
                # Adjust these selectors based on your Tableau view's structure
                self.driver.click(f'[data-filter-name="{filter_name}"]')
                
                for option in selected_options:
                    self.driver.click(f'[data-filter-option="{option}"]')
                
                # Apply filter
                self.driver.click('#filter-apply-button')
        
        # Pyppeteer method
        elif hasattr(self.driver, 'querySelector'):
            import asyncio
            
            async def apply_filters_async():
                for filter_name, selected_options in selected_filters.items():
                    await self.driver.click(f'[data-filter-name="{filter_name}"]')
                    
                    for option in selected_options:
                        await self.driver.click(f'[data-filter-option="{option}"]')
                    
                    await self.driver.click('#filter-apply-button')
            
            asyncio.get_event_loop().run_until_complete(apply_filters_async())
    
    def capture_screenshot(self, output_path):
        """
        Capture screenshot of current view
        
        :param output_path: Path to save screenshot
        :return: Path to saved screenshot
        """
        if not self.driver:
            raise ValueError("No web driver connected.")
        
        # Playwright method
        if hasattr(self.driver, 'screenshot'):
            self.driver.screenshot(path=output_path)
        
        # Pyppeteer method
        elif hasattr(self.driver, 'screenshot'):
            import asyncio
            asyncio.get_event_loop().run_until_complete(
                self.driver.screenshot({'path': output_path})
            )
        
        return output_path
    
    def generate_ppt(self, screenshots, output_ppt_path):
        """
        Generate PowerPoint with screenshots
        
        :param screenshots: List of screenshot file paths
        :param output_ppt_path: Path to save PowerPoint
        """
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        for screenshot in screenshots:
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0.5)
            pic = slide.shapes.add_picture(
                screenshot, 
                left, 
                top, 
                width=Inches(9), 
                height=Inches(6)
            )
        
        prs.save(output_ppt_path)
    
    def cleanup(self):
        """
        Close web driver and cleanup resources
        """
        if hasattr(self.driver, 'close'):
            self.driver.close()
        
        if self.playwright_context:
            self.playwright_context.stop()

def main():
    # Example usage
    tableau_url = "https://your-tableau-server.com/views/your-view"
    
    # Initialize generator
    generator = TableauFilterPPTGenerator(tableau_url)
    
    try:
        # Option 1: Use Playwright (recommended)
        try:
            page = generator.connect_with_playwright(
                username="your_username", 
                password="your_password"
            )
            
            # Extract filters
            available_filters = generator.extract_filters_with_requests()
            print("Available Filters:", available_filters)
            
            # Apply filters
            selected_filters = {
                "Region": ["North", "South"],
                "Product Category": ["Electronics"]
            }
            generator.apply_filters(selected_filters)
            
            # Capture screenshots
            screenshots = []
            for i in range(len(selected_filters)):
                screenshot_path = f"screenshot_{i}.png"
                screenshots.append(generator.capture_screenshot(screenshot_path))
            
            # Generate PowerPoint
            generator.generate_ppt(screenshots, "filtered_tableau_views.pptx")
        
        except ImportError:
            # Option 2: Use Pyppeteer if Playwright is not available
            try:
                page = generator.connect_with_pyppeteer(
                    username="your_username", 
                    password="your_password"
                )
                
                # Similar steps as Playwright method
                # ... (same code as above)
            
            except ImportError:
                # Option 3: Fallback to requests-only method (very limited)
                available_filters = generator.extract_filters_with_requests()
                print("Limited Filters Extracted:", available_filters)
                print("Full interaction not possible without web automation library")
    
    finally:
        # Always cleanup
        generator.cleanup()

if __name__ == "__main__":
    main()
